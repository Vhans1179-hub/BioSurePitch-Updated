"""
Generate a PowerPoint flowchart for the Intelligence Hub Multi-Agent System.
Similar style to the provided example.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_intelligence_hub_flowchart():
    """Create a PowerPoint presentation with the Intelligence Hub flowchart."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Add slide with blank layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 41, 59)  # Dark blue-gray
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Novartis Procurement Intelligence Hub"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add yellow line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(12.33), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(251, 191, 36)  # Yellow
    line.line.fill.background()
    
    # Helper function to create boxes
    def create_box(left, top, width, height, text, subtitle="", color=RGBColor(30, 58, 138), text_color=RGBColor(251, 191, 36)):
        """Create a styled box with text."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(59, 130, 246)  # Blue border
        shape.line.width = Pt(2)
        
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        
        # Main text
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            p2 = text_frame.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(255, 255, 255)
            p2.alignment = PP_ALIGN.CENTER
        
        return shape
    
    # Helper function to create arrows
    def create_arrow(x1, y1, x2, y2):
        """Create an arrow connector."""
        connector = slide.shapes.add_connector(
            1,  # Straight connector
            Inches(x1), Inches(y1),
            Inches(x2), Inches(y2)
        )
        connector.line.color.rgb = RGBColor(96, 165, 250)  # Light blue
        connector.line.width = Pt(3)
        return connector
    
    # Layout positions
    # Row 1: User Input
    create_box(0.5, 1.8, 2, 1, "User", "Upload Bid PDF\n+ Supplier Details")
    
    # Row 2: Orchestrator
    create_box(3.5, 1.8, 2.5, 1, "Orchestrator", "Coordinates\nAll Agents")
    
    # Row 3: Three Agents
    create_box(0.5, 3.5, 2.2, 1.2, "Technical\nAgent", "PDF Analysis\nSOP Comparison", RGBColor(30, 58, 138))
    create_box(3.2, 3.5, 2.2, 1.2, "Risk\nAgent", "Financial/ESG\nRed Flags", RGBColor(30, 58, 138))
    create_box(5.9, 3.5, 2.2, 1.2, "Financial\nAgent", "Should-Cost\nModel", RGBColor(30, 58, 138))
    
    # Row 4: Data Sources
    create_box(0.3, 5.3, 1.8, 0.7, "Gemini RAG", "PDF Extraction", RGBColor(51, 65, 85), RGBColor(255, 255, 255))
    create_box(2.3, 5.3, 1.8, 0.7, "Mock Risk DB", "Supplier Data", RGBColor(51, 65, 85), RGBColor(255, 255, 255))
    create_box(4.3, 5.3, 1.8, 0.7, "Commodity\nPrices", "Fixed Constants", RGBColor(51, 65, 85), RGBColor(255, 255, 255))
    create_box(6.3, 5.3, 1.8, 0.7, "SOP\nParameters", "Compliance Rules", RGBColor(51, 65, 85), RGBColor(255, 255, 255))
    
    # Row 5: Aggregation
    create_box(3.5, 6.5, 2.5, 0.8, "Score Aggregation", "40% Tech + 30% Risk\n+ 30% Financial")
    
    # Row 6: Output
    create_box(8.8, 1.8, 3.5, 1.2, "Executive\nSummary", "Natural Language\nRecommendation")
    create_box(8.8, 3.5, 3.5, 1.2, "Final\nRecommendation", "APPROVE/REJECT\n+ Next Steps")
    create_box(8.8, 5.3, 3.5, 0.8, "Compliance Log", "GxP Audit Trail")
    
    # Arrows - User to Orchestrator
    create_arrow(2.5, 2.3, 3.5, 2.3)
    
    # Orchestrator to Agents
    create_arrow(4.75, 2.8, 1.6, 3.5)
    create_arrow(4.75, 2.8, 4.3, 3.5)
    create_arrow(4.75, 2.8, 7.0, 3.5)
    
    # Agents to Data Sources
    create_arrow(1.3, 4.7, 1.2, 5.3)
    create_arrow(4.3, 4.7, 3.2, 5.3)
    create_arrow(7.0, 4.7, 5.2, 5.3)
    create_arrow(1.6, 4.7, 7.2, 5.3)
    
    # Agents to Aggregation
    create_arrow(1.6, 4.7, 4.0, 6.5)
    create_arrow(4.3, 4.7, 4.75, 6.5)
    create_arrow(7.0, 4.7, 5.5, 6.5)
    
    # Aggregation to Outputs
    create_arrow(6.0, 6.9, 8.8, 2.4)
    create_arrow(6.0, 6.9, 8.8, 4.1)
    create_arrow(6.0, 6.9, 8.8, 5.7)
    
    # Outputs back to User
    create_arrow(8.8, 2.4, 2.5, 2.3)
    
    # Add legend/notes at bottom
    notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(7.5), Inches(0.5))
    notes_frame = notes_box.text_frame
    notes_frame.text = "⚡ Processing Time: ~10-15 seconds  |  🔒 GxP Compliant  |  📊 Weighted Scoring: 40/30/30"
    notes_para = notes_frame.paragraphs[0]
    notes_para.font.size = Pt(12)
    notes_para.font.color.rgb = RGBColor(203, 213, 225)
    notes_para.alignment = PP_ALIGN.LEFT
    
    # Save presentation
    filename = "Novartis_Procurement_Intelligence_Hub.pptx"
    prs.save(filename)
    print(f"✅ PowerPoint created: {filename}")
    print(f"📍 Location: {filename}")
    return filename


if __name__ == "__main__":
    # Check if python-pptx is installed
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ Error: python-pptx is not installed")
        print("📦 Install it with: pip install python-pptx")
        exit(1)
    
    # Generate the flowchart
    create_intelligence_hub_flowchart()
    print("\n🎯 Flowchart shows:")
    print("   1. User uploads bid PDF with supplier details")
    print("   2. Orchestrator coordinates three specialized agents")
    print("   3. Each agent analyzes different aspects (Technical, Risk, Financial)")
    print("   4. Agents use various data sources (Gemini, Mock DBs, SOPs)")
    print("   5. Scores are aggregated with weighted formula")
    print("   6. System generates executive summary and recommendation")
    print("   7. All decisions logged for GxP compliance")